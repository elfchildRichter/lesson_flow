from __future__ import annotations

import io
import json
import logging
import math
import os
import re
import uuid
from concurrent.futures import ThreadPoolExecutor
from pypdf import PdfReader
from typing import Optional

from .models import Chunk, Deck, Document, Slide, Source
from .workflows import build_deck_graph, build_qa_graph


class DocumentStore:
    def __init__(self) -> None:
        self.documents: dict[str, Document] = {}
        self.decks: dict[str, Deck] = {}

    def add(self, document: Document) -> None:
        self.documents[document.id] = document

    def get(self, document_id: str) -> Document:
        if document_id not in self.documents:
            raise KeyError(document_id)
        return self.documents[document_id]


def _clean(text: str) -> str:
    text = re.sub(r"[\t\r]+", " ", text)
    text = re.sub(r" +", " ", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _split_page(text: str, page: int, start_index: int) -> list[Chunk]:
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    chunks: list[Chunk] = []
    buffer = ""
    for paragraph in paragraphs:
        if len(buffer) + len(paragraph) < 900:
            buffer = f"{buffer}\n{paragraph}".strip()
        else:
            if buffer:
                chunks.append(Chunk(buffer, page, start_index + len(chunks)))
            buffer = paragraph
    if buffer:
        chunks.append(Chunk(buffer, page, start_index + len(chunks)))
    if not chunks and text.strip():
        chunks.append(Chunk(text[:1200], page, start_index))
    return chunks


MATH_INDICATORS = {
    "\\int", "\\sum", "\\sqrt", "\\frac", "\\lim", "\\alpha", "\\beta", "\\gamma",
    "\\theta", "\\pi", "\\sigma", "\\infty", "\\partial", "\\matrix", "\\begin",
    "\\over", "\\vec", "\\cdot", "\\times", "\\div", "\\pm", "\\mp", "\\neq",
    "\\leq", "\\geq", "\\approx", "\\equiv", "\\propto", "\\in", "\\notin",
    "\\subset", "\\cup", "\\cap", "\\forall", "\\exists",
    "±", "∓", "≠", "≤", "≥", "≈", "≡", "∞", "∫", "∑", "∏", "√", "∛", "∜", "∂",
    "∇", "∈", "∉", "⊂", "⊆", "∪", "∩", "∧", "∨", "¬", "⇒", "⇔", "→", "↑", "↓",
    "°", "℃", "℉", "Å", "μ", "Ω", "π", "θ", "λ", "Δ", "α", "β", "γ", "δ", "ε",
    "ζ", "η", "κ", "μ", "ν", "ξ", "ρ", "σ", "τ", "φ", "χ", "ψ", "ω"
}


def _page_needs_vision(page: typing.Any, text: str) -> bool:
    """判斷 PDF 頁面是否包含圖片、掃描內容、表格結構或數學公式/符號，從而需要 VLM 多模態 Vision 解析。"""
    try:
        if hasattr(page, "get_images") and callable(page.get_images):
            if len(page.get_images()) > 0:
                return True
    except Exception:
        pass

    clean_text = text.strip()

    # 1. 若文字量極少 (例如圖片/掃描頁面)，需走 Vision 解析
    if len(clean_text) < 40:
        return True

    # 2. 檢查是否包含數學/理化公式符號
    if any(indicator in clean_text for indicator in MATH_INDICATORS):
        return True

    # 3. 檢查是否包含表格結構特徵 (如 Markdown 表格豎線 `|` 或連續 Tab 分離)
    if clean_text.count("|") >= 2 or "\t\t" in clean_text:
        return True

    return False


def parse_pdf(
    content: bytes,
    filename: str,
    ai_service: Optional[AIService] = None,
    enable_multimodal: Optional[bool] = None,
) -> Document:
    chunks: list[Chunk] = []
    total_pages = 0

    # 1. 嘗試使用 PyMuPDF + Vision (方案 A) 進行多模態直解 (含 LaTeX 與圖說)
    if enable_multimodal is None:
        enable_multimodal = os.getenv("ENABLE_MULTIMODAL_PARSING", "false").lower() in ("true", "1", "yes")
    enable_smart_routing = os.getenv("ENABLE_SMART_ROUTING", "true").lower() in ("true", "1", "yes")

    if enable_multimodal and ai_service is not None:
        try:
            import pymupdf  # PyMuPDF
            doc = pymupdf.open(stream=content, filetype="pdf")
            total_pages = len(doc)
            dpi = int(os.getenv("MULTIMODAL_DPI", "200"))
            
            max_workers = int(os.getenv("MULTIMODAL_MAX_WORKERS", "10"))
            pages_list = [(page_idx, page) for page_idx, page in enumerate(doc, 1)]

            def _process_page(item: tuple[int, pymupdf.Page]) -> tuple[int, str]:
                page_idx, page = item
                text = _clean(page.get_text() or "")

                # 頁面級智慧分流：若未開啟或判定需要 Vision，走 Vision 管道；否則直接使用本機擷取之純文字
                if enable_smart_routing and not _page_needs_vision(page, text):
                    logging.getLogger(__name__).info("第 %d 頁判定為純文字頁面，走極速本機解析管道", page_idx)
                    return page_idx, text

                try:
                    pix = page.get_pixmap(dpi=dpi)
                    img_bytes = pix.tobytes("png")
                    page_markdown = ai_service._vision_page_to_markdown(img_bytes)
                    if page_markdown and page_markdown.strip():
                        return page_idx, page_markdown
                except Exception as page_exc:
                    logging.getLogger(__name__).warning("第 %d 頁 Vision 解析失敗：%s", page_idx, page_exc)
                return page_idx, text

            actual_workers = max(1, min(max_workers, len(pages_list)))
            with ThreadPoolExecutor(max_workers=actual_workers) as executor:
                page_results = list(executor.map(_process_page, pages_list))

            page_results.sort(key=lambda x: x[0])
            for page_idx, page_markdown in page_results:
                if page_markdown and page_markdown.strip():
                    chunks.extend(_split_page(page_markdown, page_idx, len(chunks)))
        except Exception as exc:
            logging.getLogger(__name__).warning("PyMuPDF 多模態直解失敗，將自動平滑降級為 pypdf 文字模式：%s", exc)
            chunks = []

    # 2. 若多模態模式未啟動、失敗或傳回空資料，自動降級為 pypdf 文字擷取
    if not chunks:
        reader = PdfReader(io.BytesIO(content))
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception as exc:
                raise ValueError("無法讀取加密的 PDF") from exc

        total_pages = len(reader.pages)
        for page_number, page in enumerate(reader.pages, 1):
            text = _clean(page.extract_text() or "")
            chunks.extend(_split_page(text, page_number, len(chunks)))

    if not chunks:
        raise ValueError("PDF 沒有可擷取的文字與影像內容；請確認檔案正常且未毀損")

    return Document(
        id=uuid.uuid4().hex[:12],
        name=filename,
        pages=total_pages,
        chunks=chunks,
        size_bytes=len(content),
    )


def parse_text(content: str, filename: str = "AI_教案教材.md") -> Document:
    clean_text = _clean(content)
    if not clean_text:
        raise ValueError("教材內容不可為空")

    pages_list = [p.strip() for p in re.split(r"\n(?=# |\n---\n)", clean_text) if p.strip()]
    if not pages_list:
        pages_list = [clean_text]

    chunks: list[Chunk] = []
    for page_idx, page_text in enumerate(pages_list, 1):
        chunks.extend(_split_page(page_text, page_idx, len(chunks)))

    return Document(
        id=uuid.uuid4().hex[:12],
        name=filename,
        pages=len(pages_list),
        chunks=chunks,
        size_bytes=len(content.encode("utf-8")),
    )


def _prepare_openai_strict_schema(schema: dict) -> dict:
    """遞迴轉換 JSON Schema 以符合 OpenAI Strict Structured Output 的語義要求 (如包含 additionalProperties: False 與 required 欄位)。"""
    if not isinstance(schema, dict):
        return schema

    cleaned = {}
    unsupported_keys = {"minItems", "maxItems", "minimum", "maximum", "minLength", "maxLength", "pattern", "format", "default"}

    for key, value in schema.items():
        if key in unsupported_keys:
            continue
        if key == "properties" and isinstance(value, dict):
            cleaned_props = {}
            for prop_name, prop_schema in value.items():
                cleaned_props[prop_name] = _prepare_openai_strict_schema(prop_schema)
            cleaned["properties"] = cleaned_props
        elif key == "items" and isinstance(value, dict):
            cleaned["items"] = _prepare_openai_strict_schema(value)
        else:
            cleaned[key] = value

    if cleaned.get("type") == "object" or "properties" in cleaned:
        cleaned["type"] = "object"
        cleaned["additionalProperties"] = False
        if "properties" in cleaned and isinstance(cleaned["properties"], dict):
            cleaned["required"] = list(cleaned["properties"].keys())

    return cleaned


def _parse_json_response(content: str) -> dict:
    content_clean = content.strip()
    # 處理 reasoning/thinking 標籤 (包含未閉合的 <think>)
    if "</think>" in content_clean:
        content_clean = content_clean.split("</think>")[-1].strip()

    match = re.search(r"```(?:json)?\s*(\{.*\}|\[.*\])\s*```", content_clean, re.DOTALL)
    if match:
        return json.loads(match.group(1).strip())

    start_idx = content_clean.find("{")
    end_idx = content_clean.rfind("}")
    if start_idx != -1 and end_idx != -1 and end_idx > start_idx:
        json_str = content_clean[start_idx : end_idx + 1]
        try:
            return json.loads(json_str)
        except Exception:
            pass

    return json.loads(content_clean)



class AIService:
    def __init__(self) -> None:
        self.openai = None
        self.gemini_client = None
        self.ollama = None
        self._embedder = None
        self.api_key = os.getenv("OPENAI_API_KEY", "").strip()
        self.gemini_api_key = os.getenv("GEMINI_API_KEY", "").strip()
        self.qa_graph = build_qa_graph()
        self.deck_graph = build_deck_graph()

        initial_provider = os.getenv("AI_PROVIDER", "gemini").strip().lower()
        valid_providers = {"openai", "gemini", "ollama", "ollama_cloud", "ollama_local"}
        self.set_provider(initial_provider if initial_provider in valid_providers else "gemini")

    def set_provider(self, provider: str) -> dict[str, str]:
        provider = provider.strip().lower()
        if provider == "ollama":
            base_url = os.getenv("OLLAMA_BASE_URL", "")
            provider = "ollama_local" if "localhost" in base_url or "127.0.0.1" in base_url else "ollama_cloud"

        if provider not in {"openai", "gemini", "ollama_cloud", "ollama_local"}:
            raise ValueError("AI_PROVIDER 必須是 gemini, openai, ollama_cloud 或 ollama_local")

        if provider == "gemini":
            if self.gemini_client is None and self.gemini_api_key:
                try:
                    from google import genai
                    self.gemini_client = genai.Client(api_key=self.gemini_api_key)
                except Exception as exc:
                    logging.getLogger(__name__).warning("Gemini Client 初始化失敗：%s", exc)
            self.model = os.getenv("GEMINI_MODEL", "gemini-3.6-flash")
            self.embedding_model = os.getenv("GEMINI_EMBEDDING_MODEL", "gemini-embedding-2")

        elif provider == "openai":
            if not self.api_key:
                raise ValueError("未設定 OPENAI_API_KEY，無法切換至 OpenAI")
            if self.openai is None:
                from openai import OpenAI

                self.openai = OpenAI(api_key=self.api_key)
            self.model = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
            self.embedding_model = os.getenv("OPENAI_EMBEDDING_MODEL", "text-embedding-3-small")

        elif provider == "ollama_local":
            from ollama import Client
            local_url = os.getenv("OLLAMA_LOCAL_URL", "http://localhost:11434")
            local_client = Client(host=local_url)
            try:
                local_client.list()
            except Exception as exc:
                raise ValueError(f"未偵測到 Ollama 本機服務，請確認已安裝並啟動 Ollama ({local_url})。") from exc

            self.ollama = local_client
            self.model = os.getenv("OLLAMA_LOCAL_MODEL", os.getenv("OLLAMA_MODEL", "qwen3:4b"))
            self.embedding_model = os.getenv(
                "HUGGINGFACE_EMBEDDING_MODEL",
                "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2",
            )

        elif provider == "ollama_cloud":
            from ollama import Client
            cloud_url = os.getenv("OLLAMA_BASE_URL", "https://api.ollama.com")
            ollama_key = os.getenv("OLLAMA_API_KEY", "").strip()
            headers = {"Authorization": f"Bearer {ollama_key}"} if ollama_key else {}
            self.ollama = Client(host=cloud_url, headers=headers)
            self.model = os.getenv("OLLAMA_MODEL", "deepseek-v4-flash:0731")
            self.embedding_model = os.getenv("OLLAMA_EMBEDDING_MODEL", "bge-m3")

        self.provider = provider
        return self.info

    @property
    def info(self) -> dict[str, str]:
        labels = {
            "gemini": "Gemini 雲端 API",
            "openai": "OpenAI 雲端 API",
            "ollama_cloud": "Ollama 雲端 API",
            "ollama_local": "Ollama 本機服務",
        }
        label = labels.get(self.provider, "Ollama")
        return {
            "provider": self.provider,
            "provider_label": label,
            "generation_model": self.model,
            "embedding_model": self.embedding_model,
        }

    def _vision_page_to_markdown(self, img_bytes: bytes) -> str:
        """使用 VLM 進行 Vision-Native 多模態頁面直解 (提取表格、LaTeX 公式與觀念圖說)"""
        prompt = (
            "你是一個專業的教材文件 Vision-to-Markdown 解析專家。請分析這張教材頁面圖片，將其轉譯為標準 Markdown 格式：\n"
            "1. 數學公式與理化符號：請將所有單行或獨立公式轉譯為標準 LaTeX 語法（例如 `$E=mc^2$` 或 `$$\\frac{-b \\pm \\sqrt{b^2-4ac}}{2a}$$`）。\n"
            "2. 表格：請轉為標準 Markdown 表格格式。\n"
            "3. 圖片與圖表：請插入 `![圖說描述](fig)` 並詳細說明圖片呈現的觀念、數據與實驗結果。\n"
            "4. 保持原本的章節標題階層 (#, ##, ###) 與排版順序。"
        )

        import base64
        b64_str = base64.b64encode(img_bytes).decode("utf-8")

        if self.provider == "gemini" and self.gemini_client is not None:
            try:
                from google.genai import types
                response = self.gemini_client.models.generate_content(
                    model=self.model,
                    contents=[
                        types.Part.from_bytes(data=img_bytes, mime_type="image/png"),
                        prompt,
                    ],
                )
                return (response.text or "").strip()
            except Exception as exc:
                logging.getLogger(__name__).warning("Gemini Vision 解析失敗：%s", exc)

        elif self.provider == "openai" and self.openai is not None:
            try:
                response = self.openai.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": prompt},
                                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64_str}"}},
                            ],
                        }
                    ],
                    max_tokens=2000,
                )
                return (response.choices[0].message.content or "").strip()
            except Exception as exc:
                logging.getLogger(__name__).warning("OpenAI Vision 解析失敗：%s", exc)

        elif self.provider in ("ollama_cloud", "ollama_local") and self.ollama is not None:
            try:
                vision_model = os.getenv("OLLAMA_CLOUD_VISION_MODEL" if self.provider == "ollama_cloud" else "OLLAMA_LOCAL_VISION_MODEL", "qwen2-vl:7b")
                response = self.ollama.chat(
                    model=vision_model,
                    messages=[{"role": "user", "content": prompt, "images": [b64_str]}],
                )
                message = response.message if hasattr(response, "message") else response["message"]
                return (message.content if hasattr(message, "content") else message["content"]).strip()
            except Exception as exc:
                logging.getLogger(__name__).warning("Ollama Vision 解析失敗：%s", exc)

        return ""

    def _embed(self, texts: list[str], *, query: bool = False) -> list[list[float]]:
        if self.provider == "gemini" and self.gemini_client is not None:
            try:
                from google.genai import types
                vectors = []
                batch_size = 50
                for i in range(0, len(texts), batch_size):
                    batch = texts[i : i + batch_size]
                    contents = [types.Content(parts=[types.Part.from_text(text=t)]) for t in batch]
                    res = self.gemini_client.models.embed_content(
                        model=self.embedding_model,
                        contents=contents,
                    )
                    if hasattr(res, "embeddings") and res.embeddings:
                        for emb in res.embeddings:
                            vectors.append(emb.values if hasattr(emb, "values") else list(emb))
                if vectors:
                    return vectors
            except Exception as exc:
                logging.getLogger(__name__).warning("Gemini Embedding 計算失敗：%s", exc)

        if self.provider == "openai" and self.openai is not None:
            response = self.openai.embeddings.create(model=self.embedding_model, input=texts)
            return [item.embedding for item in response.data]

        if self.provider == "ollama_cloud" and self.ollama is not None:
            try:
                # 嘗試呼叫 Ollama 原生 embed API
                if hasattr(self.ollama, "embed"):
                    res = self.ollama.embed(model=self.embedding_model, input=texts)
                    embeddings = res.embeddings if hasattr(res, "embeddings") else res.get("embeddings", [])
                    if embeddings:
                        return embeddings
            except Exception as exc:
                logging.getLogger(__name__).warning("Ollama Cloud embed API 呼叫失敗，將降級為本地模型：%s", exc)

        # 延遲載入 (Lazy Loading) 本地 Hugging Face Embedding 模型 (避免在雲端模式下浪費 GB 級 RAM)
        try:
            if self._embedder is None:
                from sentence_transformers import SentenceTransformer
                hf_model = os.getenv("HUGGINGFACE_EMBEDDING_MODEL", "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2")
                self._embedder = SentenceTransformer(hf_model)
            method_name = "encode_query" if query and hasattr(self._embedder, "encode_query") else "encode_document"
            method = getattr(self._embedder, method_name, self._embedder.encode)
            vectors = method(texts, normalize_embeddings=True, show_progress_bar=False)
            return vectors.tolist() if hasattr(vectors, "tolist") else [list(vector) for vector in vectors]
        except Exception as exc:
            raise RuntimeError(f"Embedding 模型載入或推論失敗：{exc}") from exc

    def index(self, document: Document) -> None:
        document.vectors = self._embed([chunk.text for chunk in document.chunks])

    def retrieve(self, document: Document, query: str, limit: int = 4) -> list[tuple[Chunk, float]]:
        if not document.vectors:
            self.index(document)
        query_vector = self._embed([query], query=True)[0]
        if document.vectors and len(document.vectors[0]) != len(query_vector):
            self.index(document)
        ranked = []
        for chunk, vector in zip(document.chunks, document.vectors):
            dot = sum(a * b for a, b in zip(query_vector, vector))
            denom = math.sqrt(sum(a * a for a in query_vector)) * math.sqrt(sum(b * b for b in vector))
            ranked.append((chunk, dot / denom if denom else 0.0))
        ranked.sort(key=lambda pair: pair[1], reverse=True)
        return ranked[:limit]

    def _text_response(self, system: str, prompt: str) -> str:
        if self.provider == "gemini" and self.gemini_client is not None:
            try:
                from google.genai import types
                response = self.gemini_client.models.generate_content(
                    model=self.model,
                    contents=prompt,
                    config=types.GenerateContentConfig(system_instruction=system, temperature=0.1),
                )
                return (response.text or "").strip()
            except Exception as exc:
                logging.getLogger(__name__).warning("Gemini 文字生成失敗：%s", exc)

        if self.provider == "openai" and self.openai is not None:
            try:
                response = self.openai.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": system},
                        {"role": "user", "content": prompt},
                    ],
                    temperature=0.1,
                )
                return (response.choices[0].message.content or "").strip()
            except Exception as exc:
                logging.getLogger(__name__).warning("OpenAI 文字生成失敗：%s", exc)
                raise RuntimeError(f"OpenAI 生成失敗：{exc}") from exc

        try:
            response = self.ollama.chat(
                model=self.model,
                messages=[{"role": "system", "content": system}, {"role": "user", "content": prompt}],
                stream=False,
                options={"temperature": 0.1},
            )
            message = response.message if hasattr(response, "message") else response["message"]
            return (message.content if hasattr(message, "content") else message["content"]).strip()
        except Exception as exc:
            raise RuntimeError(
                f"Ollama 回應失敗；請確認服務已啟動且已執行 `ollama pull {self.model}`：{exc}"
            ) from exc

    def _structured_response(self, system: str, prompt: str, schema: dict) -> dict:
        if self.provider == "gemini" and self.gemini_client is not None:
            try:
                from google.genai import types
                response = self.gemini_client.models.generate_content(
                    model=self.model,
                    contents=prompt + f"\n\n請務必輸出符合結構的 JSON。\nSchema: {json.dumps(schema, ensure_ascii=False)}",
                    config=types.GenerateContentConfig(
                        system_instruction=system,
                        response_mime_type="application/json",
                        temperature=0,
                    ),
                )
                return _parse_json_response(response.text or "{}")
            except Exception as exc:
                logging.getLogger(__name__).warning("Gemini 結構化輸出失敗，準備降級：%s", exc)

        if self.provider == "openai" and self.openai is not None:
            try:
                strict_schema = _prepare_openai_strict_schema(schema)
                response = self.openai.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": system},
                        {"role": "user", "content": prompt},
                    ],
                    response_format={
                        "type": "json_schema",
                        "json_schema": {
                            "name": "structured_output",
                            "strict": True,
                            "schema": strict_schema,
                        },
                    },
                    temperature=0,
                )
                content = response.choices[0].message.content or "{}"
                return _parse_json_response(content)
            except Exception as exc:
                logging.getLogger(__name__).warning("OpenAI 結構化輸出 (strict json_schema) 失敗，將降級為 json_object 模式：%s", exc)
                try:
                    response = self.openai.chat.completions.create(
                        model=self.model,
                        messages=[
                            {"role": "system", "content": system},
                            {"role": "user", "content": prompt + f"\n\n請務必僅輸出合法 JSON，Schema 規定如下：\n{json.dumps(schema, ensure_ascii=False)}"},
                        ],
                        response_format={"type": "json_object"},
                        temperature=0,
                    )
                    content = response.choices[0].message.content or "{}"
                    return _parse_json_response(content)
                except Exception as fallback_exc:
                    raise RuntimeError(f"OpenAI 內容生成失敗：{fallback_exc}") from fallback_exc
        
        messages = [
            {"role": "system", "content": system},
            {"role": "user", "content": prompt + "\n\n請嚴格依照指定 JSON schema 輸出。"},
        ]
        try:
            response = self.ollama.chat(
                model=self.model,
                messages=messages,
                format=schema,
                stream=False,
                options={"temperature": 0},
            )
            message = response.message if hasattr(response, "message") else response["message"]
            content = message.content if hasattr(message, "content") else message["content"]
            return _parse_json_response(content)
        except Exception as exc_schema:
            try:
                fallback_prompt = prompt + f"\n\n請務必僅輸出合法 JSON，Schema 規定如下：\n{json.dumps(schema, ensure_ascii=False)}"
                response = self.ollama.chat(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": system},
                        {"role": "user", "content": fallback_prompt},
                    ],
                    format="json",
                    stream=False,
                    options={"temperature": 0},
                )
                message = response.message if hasattr(response, "message") else response["message"]
                content = message.content if hasattr(message, "content") else message["content"]
                return _parse_json_response(content)
            except Exception as exc_json:
                raise RuntimeError(
                    f"Ollama 結構化輸出失敗；請確認模型 `{self.model}` 可用：{exc_schema} | {exc_json}"
                ) from exc_json

    def ask(
        self, document: Document, question: str, enable_web_search: bool = False
    ) -> tuple[str, list[Source], str]:
        initial_state = {
            "question": question,
            "document": document,
            "enable_web_search": enable_web_search,
            "ai_service": self,
        }
        result = self.qa_graph.invoke(initial_state)
        answer = result.get("answer", "")
        sources = result.get("sources", [])
        return answer, sources, self.provider

    def generate_deck(
        self,
        document: Document,
        audience: str,
        tone: str,
        slide_count: int,
        duration: int,
        enable_web_search: bool = False,
        language: str = "zh-TW",
    ) -> Deck:
        initial_state = {
            "document": document,
            "audience": audience,
            "tone": tone,
            "language": language,
            "slide_count": slide_count,
            "duration": duration,
            "enable_web_search": enable_web_search,
            "ai_service": self,
        }
        result = self.deck_graph.invoke(initial_state)
        deck = result.get("deck")
        if not deck:
            raise RuntimeError("簡報生成圖未傳回有效 Deck 物件")
        return deck


def make_pptx(deck: Deck) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for index, item in enumerate(deck.slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(247, 247, 242)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
        accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor(222, 91, 55); accent.line.fill.background()
        number = slide.shapes.add_textbox(Inches(11.9), Inches(0.45), Inches(0.8), Inches(0.4))
        p = number.text_frame.paragraphs[0]; p.text = f"{index + 1:02d}"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = RGBColor(222, 91, 55); p.alignment = PP_ALIGN.RIGHT

        icon_str = getattr(item, "icon", "💡") or "💡"
        title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.65), Inches(10.8), Inches(1.1))
        p = title_box.text_frame.paragraphs[0]
        p.text = f"{icon_str}  {item.title}"
        p.font.name = "Noto Sans TC"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(27, 35, 32)

        # 左側重點內容欄位
        body = slide.shapes.add_textbox(Inches(0.85), Inches(1.95), Inches(6.6), Inches(4.8))
        tf = body.text_frame; tf.word_wrap = True
        for bullet_index, bullet in enumerate(item.bullets):
            p = tf.paragraphs[0] if bullet_index == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.name = "Noto Sans TC"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(55, 65, 61)
            p.space_after = Pt(14)

        # 右側觀念圖解視覺卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.95), Inches(4.7), Inches(4.8))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(222, 91, 55); card.line.width = Pt(1.5)

        card_title_box = slide.shapes.add_textbox(Inches(7.95), Inches(2.1), Inches(4.4), Inches(0.5))
        p = card_title_box.text_frame.paragraphs[0]
        p.text = "📐 觀念邏輯架構"
        p.font.name = "Noto Sans TC"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(222, 91, 55)

        # 提煉圖解核心步驟 (最多 3 階段)
        steps = item.bullets[:3] if item.bullets else ["核心觀念說明"]
        step_y_starts = [2.7, 4.0, 5.3] if len(steps) >= 3 else ([2.7, 4.2] if len(steps) == 2 else [3.2])
        step_heights = [1.1, 1.1, 1.1] if len(steps) >= 3 else ([1.3, 1.3] if len(steps) == 2 else [2.0])

        step_labels = ["① 觀念起點", "② 核心機制", "③ 應用成果"]
        for idx, text_item in enumerate(steps):
            y_pos = step_y_starts[idx]
            h_pos = step_heights[idx]
            step_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(y_pos), Inches(4.3), Inches(h_pos))
            step_bg.fill.solid()
            if idx == 0:
                step_bg.fill.fore_color.rgb = RGBColor(255, 245, 242)
                step_bg.line.color.rgb = RGBColor(222, 91, 55)
            elif idx == 1:
                step_bg.fill.fore_color.rgb = RGBColor(237, 242, 247)
                step_bg.line.color.rgb = RGBColor(74, 85, 104)
            else:
                step_bg.fill.fore_color.rgb = RGBColor(235, 248, 255)
                step_bg.line.color.rgb = RGBColor(43, 108, 176)
            step_bg.line.width = Pt(1.0)

            tf_step = step_bg.text_frame
            tf_step.word_wrap = True
            p_lbl = tf_step.paragraphs[0]
            p_lbl.text = step_labels[idx] if idx < len(step_labels) else f"重點 {idx+1}"
            p_lbl.font.name = "Noto Sans TC"
            p_lbl.font.size = Pt(12)
            p_lbl.font.bold = True
            p_lbl.font.color.rgb = RGBColor(222, 91, 55) if idx == 0 else (RGBColor(74, 85, 104) if idx == 1 else RGBColor(43, 108, 176))

            p_txt = tf_step.add_paragraph()
            str_item = str(text_item)
            p_txt.text = str_item[:45] + ("..." if len(str_item) > 45 else "")
            p_txt.font.name = "Noto Sans TC"
            p_txt.font.size = Pt(13)
            p_txt.font.color.rgb = RGBColor(45, 55, 72)

        notes = slide.notes_slide.notes_text_frame
        notes.text = item.speaker_notes + (f"\n\n資料來源頁碼：{', '.join(map(str, item.source_pages))}" if item.source_pages else "")
    output = io.BytesIO(); prs.save(output); return output.getvalue()


def make_script(deck: Deck) -> str:
    lines = [f"# {deck.title}", "", deck.subtitle, ""]
    for i, slide in enumerate(deck.slides, 1):
        lines.extend([f"## {i}. {slide.title}", "", slide.speaker_notes, "", f"> 教材頁碼：{', '.join(map(str, slide.source_pages)) or '—'}", ""])
    return "\n".join(lines)
